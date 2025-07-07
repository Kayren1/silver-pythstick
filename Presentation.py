from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Layouts
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Function to add slide with title, content, and an image placeholder
def add_slide_with_image(title, content, image_desc):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    text_box = slide.shapes.placeholders[1]
    text_box.text = content
    
    # Add a rectangle as an image placeholder
    left = Inches(5.5)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(3)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)  # light grey
    shape.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    
    # Add description text inside the placeholder
    text_frame = shape.text_frame
    p = text_frame.add_paragraph()
    p.text = f"Image: {image_desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = 1  # center

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "AI & Robotics"
subtitle = slide.placeholders[1]
subtitle.text = "Building Intelligent Physical Agents\nBased on Chapter 25 of Russell & Norvig’s AIMA\nPresented by: [Your Name]"

# Slide 2: Introduction to Robotics
add_slide_with_image(
    "Introduction to Robotics",
    "- Robots are physical agents manipulating the world with sensors and effectors.\n"
    "- Three main categories:\n"
    "  1. Manipulators (e.g., industrial arms)\n"
    "  2. Mobile Robots (e.g., Mars rovers)\n"
    "  3. Mobile Manipulators (e.g., Honda’s Asimo)\n"
    "- Key challenges: partial observability, stochasticity, real-time constraints.",
    "Industrial robotic arm"
)

# Slide 3: Robot Hardware (Sensors & Effectors)
add_slide_with_image(
    "Robot Hardware (Sensors & Effectors)",
    "- Sensors:\n"
    "  • Passive (cameras), Active (lidar, sonar)\n"
    "  • Range finders: Time-of-flight cameras, scanning lidar\n"
    "  • Location: GPS, IMUs, wheel odometry\n"
    "- Effectors:\n"
    "  • Manipulation: multi-DOF arms, grippers\n"
    "  • Locomotion: wheels, legs, thrusters\n"
    "  • Power: electric motors, pneumatic, hydraulic",
    "Sensors and robotic gripper"
)

# Slide 4: Robotic Perception
add_slide_with_image(
    "Robotic Perception",
    "- Estimate pose, map, object identity from sensor data.\n"
    "- Belief state via probabilistic filtering (Kalman, particle filters).\n"
    "- Localization: Kalman filters, particle filters.\n"
    "- Mapping: Occupancy grids, feature-based maps (\n   SLAM).\n"
    "- Object recognition: 3D range images, vision techniques.\n"
    "- Sensor fusion: combine camera, lidar, IMU.",
    "Lidar and camera sensor fusion"
)

# Slide 5: Planning to Move (Configuration Space)
add_slide_with_image(
    "Planning to Move (Configuration Space)",
    "- Configuration space: complete robot pose (e.g., (x,y,θ) or joint angles).\n"
    "- Motion planning:\n"
    "  • Cell decomposition: divide free C-space into cells; graph search (A*).\n"
    "  • Roadmaps: sampling-based (PRM, RRT) or analytical (visibility graph).\n"
    "- Collision avoidance: C-space obstacles by inflating workspace obstacles.",
    "Motion planning trajectory"
)

# Slide 6: Planning Uncertain Movements
add_slide_with_image(
    "Planning Uncertain Movements",
    "- Uncertainty: sensor noise, actuator variability, dynamic obstacles.\n"
    "- Decision-theoretic planning:\n"
    "  • MDPs: value/policy iteration.\n"
    "  • POMDPs: maintain belief over C-space; approximate solvers.\n"
    "- Potential-field methods: repulsive/attractive potentials; local minima.\n"
    "- Vector Field Histograms (VFH): polar histogram of obstacle densities.",
    "Robot navigating uncertain terrain"
)

# Slide 7: Moving (Control & Execution)
add_slide_with_image(
    "Moving (Control & Execution)",
    "- Control as dynamical system:\n"
    "  • Inverse kinematics: compute joint angles for end-effector pose.\n"
    "  • Inverse dynamics: compute torques for desired accelerations.\n"
    "- Trajectory tracking: smooth trajectories; feedback controllers (PID, impedance).\n"
    "- Legged locomotion:\n"
    "  • Gait generation: statically vs. dynamically stable gaits.\n"
    "  • Online foot-placement optimization.\n"
    "- Compliance & safety: spring-loaded arms; haptic feedback for grasping.",
    "Robot arm executing trajectory"
)

# Slide 8: Robotic Software Architectures
add_slide_with_image(
    "Robotic Software Architectures",
    "- Three-layer architecture:\n"
    "  1. Deliberative: global planning\n"
    "  2. Executive: sequence/manage behaviors\n"
    "  3. Reactive: fast reflexive responses\n"
    "- Subsumption architecture: layered behaviors; higher layers subsume lower.\n"
    "- Behavior-based & hybrid: concurrent behaviors (GRL, COLBERT).\n"
    "- Multi-robot/swarm: decentralized local rules; coordination (Mataric, Parker).",
    "Software architecture diagram"
)

# Slide 9: Application Domains
add_slide_with_image(
    "Application Domains",
    "1. Industrial Manipulation:\n"
    "   • Factory automation, surgical assistants.\n"
    "2. Service Robotics:\n"
    "   • Domestic (Roomba), public kiosks, healthcare delivery.\n"
    "3. Entertainment & Competitions:\n"
    "   • RoboCup, AAAI mobile robot competitions.\n"
    "4. Exploration & Defense:\n"
    "   • Mars rovers (Sojourner, Spirit/Opportunity), UAVs.\n"
    "5. Human Augmentation:\n"
    "   • Prosthetics, exoskeletons, teleoperation (ROVs).\n"
    "6. Swarm & Multi-Agent:\n"
    "   • Cooperative UGV convoys, warehouse swarms (Kiva robots).",
    "Various robots in different domains"
)

# Slide 10: Summary & Takeaways
add_slide_with_image(
    "Summary & Takeaways",
    "- Robotics = sensing + acting + reasoning.\n"
    "- Must handle continuous spaces, dynamic environments, uncertainty.\n"
    "- Probabilistic perception: Kalman & particle filters for localization/mapping.\n"
    "- Configuration-space planning: cell decomposition, roadmaps.\n"
    "- Robust control: feedback controllers, impedance control, legged algorithms.\n"
    "- Architectures drive reactivity & scalability.\n"
    "- Applications: from factories to Mars, from households to healthcare.",
    "Collage of robotics concepts"
)

# Slide 11: Q&A
slide = prs.slides.add_slide(content_slide_layout)
slide.shapes.title.text = "Q&A"
slide.placeholders[1].text = "Thank you!\nAny questions or comments?"

# Save the presentation
file_path = "/mnt/data/AI_&_Robotics.pptx"
prs.save(file_path)

file_path
